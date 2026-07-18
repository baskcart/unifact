"""Build UniFact NJ AI Hub Batch 2 pitch deck (PPTX)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "UniFact-NJ-AI-Hub-Pitch-Deck.pptx"
LOGO = Path(r"C:\Users\admin\git\dahg-ai\public\brand\unifact-logo-locked-v1.png")

BG = RGBColor(0x0A, 0x0A, 0x1A)
WHITE = RGBColor(0xF5, 0xF5, 0xF7)
MUTED = RGBColor(0xA8, 0xAD, 0xBA)
ACCENT = RGBColor(0x10, 0xB9, 0x81)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
CARD = RGBColor(0x14, 0x14, 0x28)
LINE = RGBColor(0x2A, 0x2A, 0x40)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, size=18, bold=False, color=WHITE, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set(qn("a:anchor"), {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_paragraph(tf, text, size=16, bold=False, color=WHITE, space_before=6, space_after=2, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def paint_bg(slide):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    # Move to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def accent_bar(slide, left=Inches(0.7), top=Inches(0.55), width=Inches(0.9), height=Inches(0.08)):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_logo(slide, left=Inches(0.7), top=Inches(0.4), width=Inches(0.55)):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), left, top, width=width)


def footer(slide, text="UniFact · unifact.ai · NJ AI Hub Accelerator Batch 2"):
    add_textbox(slide, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3), text, size=11, color=MUTED)


def title_block(slide, title, subtitle=None):
    accent_bar(slide)
    add_textbox(slide, Inches(0.7), Inches(0.75), Inches(12), Inches(0.55), title, size=32, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(1.35), Inches(12), Inches(0.4), subtitle, size=16, color=MUTED)


def bullet_block(slide, left, top, width, height, lines, size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=WHITE)
    return box


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    paint_bg(slide)
    return slide


def slide_title(prs):
    s = new_slide(prs)
    add_logo(s, left=Inches(0.7), top=Inches(0.55), width=Inches(0.85))
    add_textbox(s, Inches(1.75), Inches(0.7), Inches(10), Inches(0.5), "UniFact", size=22, bold=True, color=WHITE)
    add_textbox(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.9), "UniFact", size=54, bold=True, color=WHITE)
    add_textbox(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.5), "One Fact. One Truth.", size=28, color=ACCENT)
    add_textbox(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.4), "Organizational truth for AI agents", size=20, color=MUTED)
    add_textbox(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.35), "https://unifact.ai", size=18, color=VIOLET)
    card(s, Inches(0.7), Inches(5.1), Inches(7.5), Inches(1.15))
    add_textbox(s, Inches(0.95), Inches(5.25), Inches(7), Inches(0.35), "NJ AI Hub Accelerator  ·  Batch 2 Application", size=16, bold=True, color=WHITE)
    add_textbox(s, Inches(0.95), Inches(5.65), Inches(7), Inches(0.35), "MVP live  ·  Product (SaaS)  ·  West Windsor–accessible", size=14, color=MUTED)
    footer(s, "unifact.ai  ·  New Jersey AI ecosystem")


def slide_problem(prs):
    s = new_slide(prs)
    title_block(s, "The problem", "Agents are shipping. Organizational truth is not.")
    bullets = [
        "• Companies now run coding agents, ops assistants, and customer chatbots side by side.",
        "• Each is fed from different sources: READMEs, wikis, Slack, outdated PDFs, tribal knowledge.",
        "• Result: conflicting answers, repeated corrections, and actions on stale or unauthorized “facts.”",
    ]
    bullet_block(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.8), bullets, size=18)
    card(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.35))
    add_textbox(
        s,
        Inches(1.0),
        Inches(5.25),
        Inches(11.3),
        Inches(0.9),
        "AI fails organizations not only because models are wrong — because the organization has no single, governed truth layer agents can trust.",
        size=16,
        bold=True,
        color=ACCENT,
    )
    footer(s)


def slide_why_now(prs):
    s = new_slide(prs)
    title_block(s, "Why now", "Work agents entered the mainstream.")
    bullets = [
        "• Teams already use Cursor, Claude, Codex, and similar shells in daily shipping work.",
        "• Those agents lack a shared, reviewable registry of what the company has decided is true — brand, policy, infra, product rules, customer commitments.",
        "• Vector RAG and prompt stuffing retrieve text. They do not give propose → review → publish → pull with audit history.",
        "• Window: own the “org truth” layer before every vendor bolts a private memory silo onto their agent.",
    ]
    bullet_block(s, Inches(0.7), Inches(2.0), Inches(12), Inches(4.2), bullets, size=18)
    footer(s)


def slide_solution(prs):
    s = new_slide(prs)
    title_block(s, "Solution", "UniFact is the organizational fact registry.")
    add_textbox(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.4), "One place for authoritative Uni Facts — with a lifecycle agents understand:", size=16, color=MUTED)
    steps = [
        ("1", "Propose", "Builders and agents suggest facts"),
        ("2", "Review", "Curators approve or reject"),
        ("3", "Publish", "Production truth"),
        ("4", "Pull", "Agents Fact Check before they act"),
    ]
    x = Inches(0.7)
    for num, label, desc in steps:
        card(s, x, Inches(2.5), Inches(2.85), Inches(2.4))
        add_textbox(s, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.45), num, size=28, bold=True, color=ACCENT)
        add_textbox(s, x + Inches(0.2), Inches(3.3), Inches(2.4), Inches(0.4), label, size=20, bold=True, color=WHITE)
        add_textbox(s, x + Inches(0.2), Inches(3.8), Inches(2.4), Inches(0.8), desc, size=14, color=MUTED)
        x += Inches(3.1)
    add_textbox(
        s,
        Inches(0.7),
        Inches(5.3),
        Inches(12),
        Inches(0.9),
        "Not a data warehouse. Not another wiki.\nThe governance layer between humans, documents, and agents.  ·  One Fact. One Truth.",
        size=16,
        color=WHITE,
    )
    footer(s)


def slide_product(prs):
    s = new_slide(prs)
    title_block(s, "Product", "What buyers get")
    rows = [
        ("Fact registry", "Namespaced Uni Facts with versioning, supersede, retract"),
        ("Work-agent MCP", "Fact Check in Cursor / Claude / Codex-class tools"),
        ("Publish workflow", "Propose → review → publish (Git-like for truth)"),
        ("Customer agents", "Ground chatbots in published facts (white-label; default: Uni)"),
        ("Sync", "Local working store + cloud origin registry"),
    ]
    y = Inches(1.95)
    for title, desc in rows:
        card(s, Inches(0.7), y, Inches(11.9), Inches(0.72))
        add_textbox(s, Inches(0.95), y + Inches(0.12), Inches(3.2), Inches(0.45), title, size=16, bold=True, color=ACCENT)
        add_textbox(s, Inches(4.2), y + Inches(0.15), Inches(8), Inches(0.45), desc, size=15, color=WHITE)
        y += Inches(0.82)
    add_textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.35), "Live: https://unifact.ai  ·  Self-serve: GitHub OAuth to create an org registry", size=14, color=MUTED)
    footer(s)


def slide_traction(prs):
    s = new_slide(prs)
    title_block(s, "Demo / traction", "Proof today (MVP)")
    cards = [
        ("Live product", "unifact.ai is up — registry, MCP, publish workflow"),
        ("Dogfooding", "UniFact org runs on UniFact; agents Fact Check before org-dependent work"),
        ("Ready for pilots", "Seeking design partners & enterprise pilots — why the Hub fits"),
    ]
    x = Inches(0.7)
    for title, body in cards:
        card(s, x, Inches(2.1), Inches(3.85), Inches(3.2))
        add_textbox(s, x + Inches(0.25), Inches(2.4), Inches(3.3), Inches(0.5), title, size=18, bold=True, color=ACCENT)
        add_textbox(s, x + Inches(0.25), Inches(3.1), Inches(3.3), Inches(1.8), body, size=15, color=WHITE)
        x += Inches(4.05)
    add_textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.6), "Optional: attach a 20–30s Fact Check → publish → agent behavior clip.", size=14, color=MUTED)
    footer(s)


def slide_market(prs):
    s = new_slide(prs)
    title_block(s, "Market fit", "Aligned to Hub / Plug and Play focus areas")
    items = [
        ("Future of Work", "Shared truth for employee/work agents so AI labor is reliable"),
        ("Big Data + AI", "Trusted context layer for agentic systems — authority, not bulk storage"),
        ("Customer Engagement", "Customer-facing agents grounded in published Uni Facts only"),
    ]
    y = Inches(2.0)
    for title, body in items:
        card(s, Inches(0.7), y, Inches(11.9), Inches(1.05))
        add_textbox(s, Inches(0.95), y + Inches(0.18), Inches(11.4), Inches(0.35), title, size=17, bold=True, color=ACCENT)
        add_textbox(s, Inches(0.95), y + Inches(0.55), Inches(11.4), Inches(0.4), body, size=15, color=WHITE)
        y += Inches(1.2)
    add_textbox(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.7), "NJ angle: pharma, finance, telecom, and public-sector adjacency need auditable agent context — propose/review/publish beats “stuff the prompt.”", size=14, color=MUTED)
    footer(s)


def slide_buyers(prs):
    s = new_slide(prs)
    title_block(s, "Who buys", "Primary ICP — first 12 months")
    add_textbox(s, Inches(0.7), Inches(1.95), Inches(12), Inches(0.55), "Teams already deploying AI work agents who keep re-explaining brand, infra policy, product decisions, and support rules.", size=16, color=MUTED)
    personas = [
        ("Founder / CTO", "AI-adopting SMB or product org"),
        ("Platform / DevEx", "Wiring agents into the SDLC"),
        ("Ops / Knowledge", "Tired of agents inventing policy"),
    ]
    x = Inches(0.7)
    for title, body in personas:
        card(s, x, Inches(2.8), Inches(3.85), Inches(1.8))
        add_textbox(s, x + Inches(0.25), Inches(3.05), Inches(3.3), Inches(0.4), title, size=17, bold=True, color=WHITE)
        add_textbox(s, x + Inches(0.25), Inches(3.55), Inches(3.3), Inches(0.7), body, size=14, color=MUTED)
        x += Inches(4.05)
    card(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.2))
    add_textbox(s, Inches(0.95), Inches(5.2), Inches(11.4), Inches(0.35), "Land → Expand", size=16, bold=True, color=ACCENT)
    add_textbox(s, Inches(0.95), Inches(5.6), Inches(11.4), Inches(0.4), "Land: one org registry + MCP in daily agent work.  Expand: more agents, namespaces, curator seats, customer-agent deployments.", size=14, color=WHITE)
    footer(s)


def slide_model(prs):
    s = new_slide(prs)
    title_block(s, "Business model", "Product company. SaaS first.")
    rows = [
        ("Core", "Org / seat / usage subscription for the registry"),
        ("Onboarding (optional)", "Fixed-scope “Agent Truth Launch” — seed Uni Facts, wire MCP, set curator workflow"),
        ("Not this", "Open-ended AI consulting firm"),
    ]
    y = Inches(2.0)
    for title, body in rows:
        card(s, Inches(0.7), y, Inches(11.9), Inches(1.05))
        add_textbox(s, Inches(0.95), y + Inches(0.2), Inches(11.4), Inches(0.35), title, size=17, bold=True, color=ACCENT if title != "Not this" else VIOLET)
        add_textbox(s, Inches(0.95), y + Inches(0.55), Inches(11.4), Inches(0.4), body, size=15, color=WHITE)
        y += Inches(1.2)
    add_textbox(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.55), "Implementation packages fund early learning and design partners. Narrative and roadmap stay product-led. Pricing in validation.", size=14, color=MUTED)
    footer(s)


def slide_why_hub(prs):
    s = new_slide(prs)
    title_block(s, "Why NJ AI Hub", "What we want from the accelerator")
    asks = [
        "Enterprise & university pilots through Hub / Plug and Play network",
        "Mentorship on commercialization and ICP sharpening",
        "Visibility in NJ’s AI ecosystem (West Windsor coworking during program)",
        "Pathways to corporates who need trustworthy agent operations",
    ]
    bullet_block(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.8), [f"• {a}" for a in asks], size=17)
    card(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.3))
    add_textbox(s, Inches(0.95), Inches(5.2), Inches(11.4), Inches(0.9), "Zero equity / no tied funding is a fit: we need distribution and pilots, not a pitched valuation.\nIn return: a live MVP in Future of Work / agent infrastructure, ready for design-partner pilots in Batch 2 (Sep–Dec 2026).", size=14, color=WHITE)
    footer(s)


def slide_roadmap(prs):
    s = new_slide(prs)
    title_block(s, "Roadmap", "Next 12 months")
    phases = [
        ("Now → 90 days", "Close design partners; tighten onboarding; polish self-serve org creation"),
        ("6 months", "Multi-org patterns; stronger curator UX; pilot case studies"),
        ("12 months", "Enterprise readiness (SSO beyond GitHub, compliance narratives); scale seats"),
    ]
    x = Inches(0.7)
    for title, body in phases:
        card(s, x, Inches(2.2), Inches(3.85), Inches(3.0))
        add_textbox(s, x + Inches(0.25), Inches(2.5), Inches(3.3), Inches(0.55), title, size=18, bold=True, color=ACCENT)
        add_textbox(s, x + Inches(0.25), Inches(3.3), Inches(3.3), Inches(1.5), body, size=15, color=WHITE)
        x += Inches(4.05)
    add_textbox(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.7), "Foundation in place: local SQLite working store + PostgreSQL origin registry; cloud-neutral upstream sync.", size=14, color=MUTED)
    footer(s)


def slide_team_ask(prs):
    s = new_slide(prs)
    title_block(s, "Team + ask", "Admit UniFact to NJ AI Hub Accelerator Batch 2")
    card(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(2.6))
    add_textbox(s, Inches(0.95), Inches(2.15), Inches(5.2), Inches(0.4), "Team", size=16, bold=True, color=ACCENT)
    add_textbox(s, Inches(0.95), Inches(2.7), Inches(5.2), Inches(1.5), "[FOUNDER NAME]\nFounder · product & platform\n\n[Optional co-founder / advisor]\n[ROLE]", size=15, color=WHITE)

    card(s, Inches(6.7), Inches(1.95), Inches(5.9), Inches(2.6))
    add_textbox(s, Inches(6.95), Inches(2.15), Inches(5.4), Inches(0.4), "The ask", size=16, bold=True, color=ACCENT)
    add_textbox(
        s,
        Inches(6.95),
        Inches(2.7),
        Inches(5.4),
        Inches(1.6),
        "1. Win 2–3 design-partner pilots\n2. Mentorship + corporate intros for agent-ops / Future of Work buyers\n3. Ship commercialization milestones by Dec 2026 Expo",
        size=14,
        color=WHITE,
    )

    card(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.35))
    add_textbox(s, Inches(0.95), Inches(5.1), Inches(11.4), Inches(0.35), "Contact", size=15, bold=True, color=ACCENT)
    add_textbox(s, Inches(0.95), Inches(5.5), Inches(11.4), Inches(0.4), "[EMAIL]  ·  https://unifact.ai", size=18, bold=True, color=WHITE)
    footer(s)


def slide_competition(prs):
    s = new_slide(prs)
    title_block(s, "Appendix — competitive positioning", "Authority beats retrieval")
    rows = [
        ("Wikis / Notion / Confluence", "Human docs; agents don’t Fact Check a governed publish layer"),
        ("Vector RAG / “memory”", "Retrieval ≠ authority; weak review, supersede, audit"),
        ("Custom prompt stuffing", "Fragile, unshared across agents"),
        ("UniFact", "Registry + lifecycle + MCP pull for universal agents"),
    ]
    y = Inches(1.95)
    for i, (title, body) in enumerate(rows):
        card(s, Inches(0.7), y, Inches(11.9), Inches(0.9))
        color = ACCENT if i == 3 else WHITE
        add_textbox(s, Inches(0.95), y + Inches(0.12), Inches(4.0), Inches(0.55), title, size=15, bold=True, color=color)
        add_textbox(s, Inches(5.1), y + Inches(0.18), Inches(7.2), Inches(0.55), body, size=14, color=MUTED if i < 3 else WHITE)
        y += Inches(1.0)
    footer(s)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_why_now(prs)
    slide_solution(prs)
    slide_product(prs)
    slide_traction(prs)
    slide_market(prs)
    slide_buyers(prs)
    slide_model(prs)
    slide_why_hub(prs)
    slide_roadmap(prs)
    slide_team_ask(prs)
    slide_competition(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
